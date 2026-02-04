#!/usr/bin/env python3
"""
STEMI EMS Presentation - Enhanced Visual Design
With EKG images, flowcharts, and memorable mnemonics
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import os

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============== ENHANCED COLOR SCHEME ==============
# Primary palette - Modern medical aesthetic
BG_DARK = RGBColor(13, 17, 23)        # Near-black for main bg
BG_CARD = RGBColor(22, 27, 34)        # Dark card background
ACCENT_RED = RGBColor(239, 68, 68)    # Urgent/Alert red
ACCENT_BLUE = RGBColor(59, 130, 246)  # Info blue
ACCENT_GREEN = RGBColor(34, 197, 94)  # Success green
ACCENT_YELLOW = RGBColor(250, 204, 21) # Warning yellow
ACCENT_PURPLE = RGBColor(168, 85, 247) # Purple for tips
TEXT_WHITE = RGBColor(248, 250, 252)  # Pure white text
TEXT_MUTED = RGBColor(148, 163, 184)  # Muted gray
TEXT_HIGHLIGHT = RGBColor(255, 255, 255)

# Image directory
IMG_DIR = "/Users/vi/.openclaw/workspace/projects/stemi-presentation/images"

def add_background(slide, color=BG_DARK):
    """Add solid background to slide"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_title_bar(slide, title, accent_color=ACCENT_RED):
    """Add styled title bar with accent stripe"""
    # Accent stripe on left
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), Inches(1.2))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent_color
    stripe.line.fill.background()
    
    # Title background
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.15), 0, Inches(13.183), Inches(1.2))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = BG_CARD
    title_bg.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

def add_image_with_border(slide, img_path, left, top, width, border_color=ACCENT_BLUE):
    """Add image with colored border frame"""
    if os.path.exists(img_path):
        # Add border rectangle first (slightly larger)
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left - Inches(0.05),
            top - Inches(0.05),
            width + Inches(0.1),
            Inches(3.5) + Inches(0.1)
        )
        border.fill.solid()
        border.fill.fore_color.rgb = border_color
        border.line.fill.background()
        
        # Add image
        slide.shapes.add_picture(img_path, left, top, width)
        return True
    return False

def add_bullet_content(slide, items, top=Inches(1.5), left=Inches(0.5), width=Inches(6)):
    """Add styled bullet point content"""
    content_box = slide.shapes.add_textbox(left, top, width, Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        
        # Detect item type
        if isinstance(item, tuple):
            text, style = item
        else:
            text = item
            style = "normal"
        
        # Handle headers
        if style == "header" or (text.isupper() and len(text) < 60):
            p.text = text
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = ACCENT_YELLOW
            p.space_before = Pt(16)
        elif style == "subitem" or text.startswith("  "):
            p.text = "    → " + text.strip()
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_MUTED
            p.space_before = Pt(4)
        elif style == "warning":
            p.text = "⚠️ " + text
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = ACCENT_RED
            p.space_before = Pt(12)
        elif style == "success":
            p.text = "✓ " + text
            p.font.size = Pt(20)
            p.font.color.rgb = ACCENT_GREEN
            p.space_before = Pt(8)
        elif style == "fail":
            p.text = "✗ " + text
            p.font.size = Pt(20)
            p.font.color.rgb = ACCENT_RED
            p.space_before = Pt(8)
        elif text.strip():
            p.text = "• " + text
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_WHITE
            p.space_before = Pt(8)
        else:
            p.text = ""
            p.space_before = Pt(8)
    
    return content_box

def add_card(slide, left, top, width, height, title, content_lines, accent=ACCENT_BLUE):
    """Add a styled card with title and content"""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    card.line.color.rgb = accent
    card.line.width = Pt(2)
    
    # Card title
    title_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = accent
    
    # Card content
    content_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.5), width - Inches(0.3), height - Inches(0.6))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    first = True
    for line in content_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(4)

def create_title_slide(title, subtitle_lines, notes=""):
    """Create a dramatic title slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide)
    
    # Large heart emoji/icon area
    heart_box = slide.shapes.add_textbox(Inches(5.5), Inches(0.5), Inches(2.333), Inches(1.5))
    tf = heart_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🫀"
    p.font.size = Pt(72)
    p.alignment = PP_ALIGN.CENTER
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle lines
    for i, line in enumerate(subtitle_lines):
        p = tf.add_paragraph()
        p.text = line
        if i == 0:
            p.font.size = Pt(28)
            p.font.color.rgb = ACCENT_RED
            p.font.italic = True
        else:
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(12)
    
    # Bottom accent line
    line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(6.8), Inches(5.333), Inches(0.05))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = ACCENT_RED
    line_shape.line.fill.background()
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    
    return slide

def create_content_slide(title, accent=ACCENT_RED):
    """Create a basic content slide with title bar"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide)
    add_title_bar(slide, title, accent)
    return slide

def create_two_column_slide(title, left_content, right_content, accent=ACCENT_RED):
    """Create a two-column layout slide"""
    slide = create_content_slide(title, accent)
    add_bullet_content(slide, left_content, Inches(1.4), Inches(0.3), Inches(6))
    add_bullet_content(slide, right_content, Inches(1.4), Inches(6.8), Inches(6))
    return slide

def create_image_slide(title, img_filename, content, img_on_right=True, accent=ACCENT_RED):
    """Create slide with image and content"""
    slide = create_content_slide(title, accent)
    
    img_path = os.path.join(IMG_DIR, img_filename)
    
    if img_on_right:
        add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(5.8))
        if os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, Inches(6.5), Inches(1.5), width=Inches(6.5))
            except:
                # Fallback text if image fails
                pass
    else:
        if os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.5), width=Inches(6.5))
            except:
                pass
        add_bullet_content(slide, content, Inches(1.4), Inches(7), Inches(5.8))
    
    return slide

# ============== BUILD SLIDES ==============

# SLIDE 1: Title
create_title_slide(
    "STEMI Recognition & Management",
    [
        "Saving Hearts, Saving Time",
        "",
        "Dr. Peter Bleszynski, MD, FACC, FSCAI",
        "Director of Cardiology, North Vista Hospital",
        "March 2026"
    ],
    "Welcome everyone. Today we'll cover STEMI - one of the most time-critical emergencies you'll face."
)

# SLIDE 2: Learning Objectives
slide = create_content_slide("Learning Objectives", ACCENT_BLUE)
objectives = [
    ("AFTER THIS SESSION YOU WILL:", "header"),
    "Recognize classic and atypical STEMI patterns on 12-lead EKG",
    "Identify common STEMI mimics to avoid false activations",
    "Apply current AHA/ACC guideline-directed prehospital management",
    "Understand door-to-balloon time importance",
    "Describe key medications in ACS management"
]
add_bullet_content(slide, objectives, Inches(1.5))

# SLIDE 3: Why STEMI Matters - TIME = MUSCLE (Visual Impact)
slide = create_content_slide("⏱️ TIME = MUSCLE", ACCENT_RED)

# Big stat boxes
add_card(slide, Inches(0.3), Inches(1.5), Inches(4), Inches(2), 
         "EVERY 30 MIN DELAY", ["↑ 7.5% Mortality", "That's someone's family member"], ACCENT_RED)
add_card(slide, Inches(4.5), Inches(1.5), Inches(4), Inches(2),
         "1 BILLION CELLS/HOUR", ["Heart muscle dying", "During complete occlusion"], ACCENT_RED)
add_card(slide, Inches(8.7), Inches(1.5), Inches(4.3), Inches(2),
         "AFTER 6 HOURS", ["Most damage irreversible", "Window is closing fast"], ACCENT_RED)

# Goals box at bottom
goals_box = slide.shapes.add_textbox(Inches(0.3), Inches(4), Inches(12.7), Inches(3))
tf = goals_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🎯 THE GOLD STANDARD GOALS"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = ACCENT_YELLOW

goals = [
    ("First Medical Contact → EKG", "≤10 minutes"),
    ("Door-to-Balloon (D2B)", "≤90 minutes"),
    ("First Medical Contact → Device", "≤90 minutes")
]
for label, time in goals:
    p = tf.add_paragraph()
    p.text = f"    {label}: {time}"
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(8)

# SLIDE 4: Chain of Survival (Timeline Infographic)
slide = create_content_slide("The Chain of Survival", ACCENT_BLUE)

# Timeline boxes
stages = [
    ("🚨", "Symptom\nOnset"),
    ("📞", "911\nCall"),
    ("🚑", "EMS\nArrival"),
    ("📊", "12-Lead\nEKG"),
    ("📱", "Cath Lab\nActivation"),
    ("🏥", "Hospital"),
    ("🫀", "Wire\nCrossing")
]

box_width = 1.6
start_x = 0.5
for i, (emoji, label) in enumerate(stages):
    x = start_x + (i * (box_width + 0.25))
    
    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(box_width), Inches(1.5))
    box.fill.solid()
    if i in [2, 3, 4]:  # EMS control points
        box.fill.fore_color.rgb = ACCENT_BLUE
    else:
        box.fill.fore_color.rgb = BG_CARD
    box.line.color.rgb = ACCENT_BLUE
    box.line.width = Pt(2)
    
    # Emoji
    emoji_box = slide.shapes.add_textbox(Inches(x), Inches(1.85), Inches(box_width), Inches(0.6))
    tf = emoji_box.text_frame
    p = tf.paragraphs[0]
    p.text = emoji
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER
    
    # Label
    label_box = slide.shapes.add_textbox(Inches(x), Inches(2.5), Inches(box_width), Inches(0.8))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Arrow
    if i < len(stages) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + box_width), Inches(2.3), Inches(0.25), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_YELLOW
        arrow.line.fill.background()

# EMS impact callout
impact_box = slide.shapes.add_textbox(Inches(0.3), Inches(3.8), Inches(12.7), Inches(3.2))
tf = impact_box.text_frame
p = tf.paragraphs[0]
p.text = "🎯 WHERE EMS MAKES THE BIGGEST IMPACT"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN
for item in ["Rapid 12-lead acquisition", "Early cath lab activation (field activation)", 
             "Bypass non-PCI facilities when appropriate", "Minimize scene time",
             "Field activation saves 15-20 minutes = LIVES SAVED"]:
    p = tf.add_paragraph()
    p.text = "  ✓ " + item
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_WHITE if "saves" not in item else ACCENT_YELLOW
    p.font.bold = "saves" in item
    p.space_before = Pt(6)

# SLIDE 5: Coronary Anatomy
slide = create_content_slide("Coronary Anatomy: The Big Three", ACCENT_RED)

# Three artery cards
arteries = [
    ("🔴 LAD", "Left Anterior Descending", ["'The Widow Maker'", "Supplies: Anterior wall,", "septum, apex", "", "STEMI: V1-V4"], ACCENT_RED),
    ("🟡 LCx", "Left Circumflex", ["'The Sneaky One'", "Supplies: Lateral wall", "", "STEMI: I, aVL, V5-V6"], ACCENT_YELLOW),
    ("🟢 RCA", "Right Coronary Artery", ["Supplies: Inferior wall,", "RV, AV node", "", "STEMI: II, III, aVF", "⚠️ Watch for bradycardia!"], ACCENT_GREEN),
]

for i, (emoji, name, details, color) in enumerate(arteries):
    x = 0.3 + (i * 4.3)
    add_card(slide, Inches(x), Inches(1.5), Inches(4), Inches(3.8), f"{emoji} {name}", details, color)

# Add coronary anatomy image if available
img_path = os.path.join(IMG_DIR, "coronary-anatomy2.jpg")
if os.path.exists(img_path):
    try:
        slide.shapes.add_picture(img_path, Inches(10), Inches(5.5), width=Inches(3))
    except:
        pass

# Mnemonic at bottom
mnem_box = slide.shapes.add_textbox(Inches(0.3), Inches(5.5), Inches(9), Inches(1.5))
tf = mnem_box.text_frame
p = tf.paragraphs[0]
p.text = "💡 MNEMONIC: \"LAD in front, RCA below, LCx to the side we go\""
p.font.size = Pt(20)
p.font.italic = True
p.font.color.rgb = ACCENT_PURPLE

# SLIDE 6: STEMI Pathophysiology
slide = create_content_slide("What Happens During a STEMI", ACCENT_RED)
content = [
    ("THE CASCADE:", "header"),
    "1. Atherosclerotic plaque ruptures",
    "2. Platelet adhesion and aggregation",
    "3. Thrombus forms → Complete occlusion",
    "4. Myocardium becomes ischemic (minutes)",
    "5. Injury pattern on EKG (minutes-hours)",
    "6. Necrosis/Infarction (hours)",
    "",
    ("SALVAGEABLE MYOCARDIUM", "header"),
    "  The 'ischemic penumbra' - viable but at risk",
    "  This is what we're racing to save",
    ("  Best outcomes within 120 minutes!", "warning")
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(12.5))

# SLIDE 7: ST-Elevation Criteria
slide = create_content_slide("EKG: ST-Elevation Criteria", ACCENT_BLUE)
content = [
    ("DEFINITION OF STEMI:", "header"),
    "  Limb leads (I, II, III, aVL, aVF): ≥1 mm elevation",
    "  Precordial leads (V1-V6): ≥2 mm elevation",
    "  Must be in ≥2 contiguous (adjacent) leads",
    "",
    ("WHAT TO LOOK FOR:", "header"),
    "  ST segment elevated above baseline",
    "  'Tombstone' or convex morphology = BAD",
    "  Reciprocal changes (ST depression in opposite leads)"
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(6.5))

# Visual representation
ekg_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(6), Inches(4))
tf = ekg_box.text_frame
p = tf.paragraphs[0]
p.text = "NORMAL vs STEMI"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_YELLOW

p = tf.add_paragraph()
p.text = "Normal:  ___/\\___"
p.font.size = Pt(24)
p.font.name = "Courier New"
p.font.color.rgb = ACCENT_GREEN
p.space_before = Pt(20)

p = tf.add_paragraph()
p.text = "(ST at baseline)"
p.font.size = Pt(16)
p.font.color.rgb = TEXT_MUTED

p = tf.add_paragraph()
p.text = "STEMI:  ___/‾‾\\___"
p.font.size = Pt(24)
p.font.name = "Courier New"
p.font.color.rgb = ACCENT_RED
p.space_before = Pt(20)

p = tf.add_paragraph()
p.text = "(ST elevated - 'tombstone')"
p.font.size = Pt(16)
p.font.color.rgb = TEXT_MUTED

# SLIDE 8: Anterior STEMI with EKG image
slide = create_content_slide("Anterior STEMI (LAD Occlusion)", ACCENT_RED)
content = [
    ("THE 'WIDOW MAKER'", "header"),
    "",
    ("EKG PATTERN:", "header"),
    "  ST elevation in V1, V2, V3, V4",
    "  May extend to V5-V6 if proximal",
    "  Reciprocal depression in II, III, aVF",
    "",
    ("CLINICAL SIGNIFICANCE:", "header"),
    "  Largest territory at risk",
    "  Highest mortality of all STEMIs",
    ("  Can rapidly develop cardiogenic shock", "warning"),
    ("  Call the cath lab NOW!", "warning")
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(5.5))

# Add anterior STEMI image
img_path = os.path.join(IMG_DIR, "ECG-Extensive-anterior-MI-tombstoning-pattern.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(6.3), Inches(1.5), width=Inches(6.8))

# SLIDE 9: Inferior STEMI with EKG image  
slide = create_content_slide("Inferior STEMI (RCA/LCx)", ACCENT_YELLOW)
content = [
    ("EKG PATTERN:", "header"),
    "  ST elevation in II, III, aVF",
    "  III > II suggests RCA",
    "  Reciprocal depression in I, aVL",
    "",
    ("⚠️ CRITICAL: CHECK V4R!", "warning"),
    "  ST elevation ≥1mm = RV infarct",
    "  Changes management!",
    "",
    ("WHY IT MATTERS:", "header"),
    "  RCA feeds AV node",
    "  Watch for bradycardia/heart block"
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(5.5))

img_path = os.path.join(IMG_DIR, "ECG-Inferior-STEMI-Hyperacute-1.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(6.3), Inches(1.5), width=Inches(6.8))

# SLIDE 10: RV Infarction - CRITICAL
slide = create_content_slide("⚠️ Right Ventricular Infarction", ACCENT_RED)

# Warning banner
banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.4), Inches(12.7), Inches(0.8))
banner.fill.solid()
banner.fill.fore_color.rgb = ACCENT_RED
banner.line.fill.background()
banner_text = slide.shapes.add_textbox(Inches(0.3), Inches(1.5), Inches(12.7), Inches(0.6))
tf = banner_text.text_frame
p = tf.paragraphs[0]
p.text = "NITRO CAN KILL THESE PATIENTS - They are PRELOAD DEPENDENT!"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.alignment = PP_ALIGN.CENTER

content = [
    ("WHEN TO SUSPECT:", "header"),
    "  Inferior STEMI on EKG",
    "  Hypotension (especially after nitrates!)",
    "  JVD + clear lungs (paradox)",
    "",
    ("V4R IS YOUR FRIEND:", "header"),
    "  ≥1 mm ST elevation = RV involvement"
]
add_bullet_content(slide, content, Inches(2.4), Inches(0.3), Inches(6))

# DO/DON'T boxes
add_card(slide, Inches(6.5), Inches(2.4), Inches(3), Inches(2.2), "❌ AVOID", 
         ["Nitrates", "Morphine", "Diuretics"], ACCENT_RED)
add_card(slide, Inches(9.8), Inches(2.4), Inches(3.2), Inches(2.2), "✓ GIVE",
         ["Aggressive IV fluids", "(Preload dependent!)"], ACCENT_GREEN)

# Add V4R image
img_path = os.path.join(IMG_DIR, "ECG-Right-ventricular-infarction-1b-V4R.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.3), Inches(4.8), width=Inches(6))

img_path2 = os.path.join(IMG_DIR, "V4R-ECG-lead-placement.jpg")
if os.path.exists(img_path2):
    slide.shapes.add_picture(img_path2, Inches(6.5), Inches(4.8), width=Inches(3))

# SLIDE 11: Lateral STEMI
slide = create_content_slide("Lateral STEMI (LCx Occlusion)", ACCENT_YELLOW)
content = [
    ("THE 'SNEAKY' STEMI", "header"),
    "",
    ("EKG PATTERN:", "header"),
    "  ST elevation in I, aVL, V5, V6",
    "  Often subtle - easy to miss!",
    "  Reciprocal depression in inferior leads",
    "",
    ("CLINICAL PEARL:", "header"),
    "  Look carefully at leads I and aVL",
    "  Even 1-2mm elevation is significant",
    "  LCx wraps around the back"
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(5.5))

img_path = os.path.join(IMG_DIR, "ECG-Lateral-STEMI-1st-diagonal.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(6.3), Inches(1.5), width=Inches(6.8))

# SLIDE 12: Posterior STEMI
slide = create_content_slide("Posterior STEMI: The Hidden STEMI", ACCENT_PURPLE)
content = [
    ("STANDARD 12-LEAD SHOWS MIRROR IMAGE:", "header"),
    "  ST DEPRESSION in V1-V3 (not elevation!)",
    "  Tall, broad R waves in V1-V2",
    "  Upright T waves in V1-V3",
    "",
    ("CONFIRM WITH POSTERIOR LEADS:", "header"),
    "  V7, V8, V9 placement",
    "  ST elevation ≥0.5mm confirms it",
    "",
    ("OFTEN OCCURS WITH:", "header"),
    "  Inferior STEMI (RCA)", 
    "  Lateral STEMI (LCx)"
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(5.5))

img_path = os.path.join(IMG_DIR, "ECG-Posterior-AMI-1.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(6.3), Inches(1.5), width=Inches(6.8))

# SLIDE 13: STEMI Mimics Overview
slide = create_content_slide("STEMI Mimics: Know the Imposters", ACCENT_YELLOW)

mimics = [
    ("LVH", "Left Ventricular\nHypertrophy"),
    ("LBBB", "Left Bundle\nBranch Block"),
    ("BER", "Early\nRepolarization"),
    ("Pericarditis", "Inflammation"),
    ("Brugada", "Genetic\nChannelopathy"),
    ("Hyperkalemia", "High\nPotassium"),
]

for i, (abbrev, name) in enumerate(mimics):
    row = i // 3
    col = i % 3
    x = 0.3 + (col * 4.3)
    y = 1.5 + (row * 2.2)
    add_card(slide, Inches(x), Inches(y), Inches(4), Inches(1.8), abbrev, [name], ACCENT_YELLOW)

# Bottom reminder
reminder = slide.shapes.add_textbox(Inches(0.3), Inches(5.9), Inches(12.7), Inches(1.2))
tf = reminder.text_frame
p = tf.paragraphs[0]
p.text = "⚠️ False activations waste resources, BUT: When in doubt, ACTIVATE!"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_RED
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "Better a false positive than a missed STEMI"
p.font.size = Pt(18)
p.font.color.rgb = TEXT_MUTED
p.alignment = PP_ALIGN.CENTER

# SLIDE 14: LVH Mimic
slide = create_content_slide("STEMI Mimic: LVH", ACCENT_YELLOW)
content = [
    ("WHY IT MIMICS STEMI:", "header"),
    "  Large QRS voltage",
    "  'Strain pattern' - ST depression in V5-V6",
    "  Can have ST elevation in V1-V3",
    "",
    ("HOW TO DIFFERENTIATE:", "header"),
    "  High voltage criteria present",
    "  (S in V1 + R in V5/V6 > 35mm)",
    "  ST changes are DISCORDANT",
    "  (opposite direction to QRS)",
    "  Patient has longstanding HTN"
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(5.5))

img_path = os.path.join(IMG_DIR, "ECG-LVH-ST-elevation-not-MI.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(6.3), Inches(1.5), width=Inches(6.8))

# SLIDE 15: LBBB Mimic with Sgarbossa
slide = create_content_slide("STEMI Mimic: LBBB", ACCENT_YELLOW)
content = [
    ("THE PROBLEM:", "header"),
    "  LBBB causes secondary ST changes",
    "  Old teaching: 'New LBBB = STEMI'",
    ("  This is OUTDATED!", "warning"),
    "",
    ("MODIFIED SGARBOSSA CRITERIA:", "header"),
    "  Concordant ST elevation ≥1mm → 5 pts",
    "  ST depression ≥1mm in V1-V3 → 3 pts",
    "  Discordant ST elev ≥25% of S → 2 pts",
    "",
    ("≥3 POINTS SUGGESTS STEMI", "success"),
    "",
    "KEY: CONCORDANT = ST goes SAME",
    "direction as QRS (abnormal in LBBB)"
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(5.8))

img_path = os.path.join(IMG_DIR, "ECG-Left-Bundle-Branch-Block-LBBB-2-2.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(6.5), Inches(1.5), width=Inches(6.5))

# SLIDE 16: Early Repolarization
slide = create_content_slide("STEMI Mimic: Early Repolarization", ACCENT_YELLOW)
content = [
    ("CLASSIC FEATURES:", "header"),
    "  ST elevation in precordial leads",
    "  'Smiley face' concave morphology",
    "  'Fish hook' J-point notching",
    "  Young, healthy patients",
    "",
    ("THE KEY DIFFERENCE:", "header"),
    ("  STEMI = Convex (frowny face) ☹️", "warning"),
    ("  BER = Concave (smiley face) 🙂", "success"),
    "",
    "  No reciprocal changes in BER",
    "  BER is stable (doesn't evolve)"
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(5.5))

img_path = os.path.join(IMG_DIR, "ECG-Benign-Early-Repolarisation-BER-1.jpeg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(6.3), Inches(1.5), width=Inches(6.8))

# SLIDE 17: Pericarditis
slide = create_content_slide("STEMI Mimic: Pericarditis", ACCENT_YELLOW)
content = [
    ("EKG FEATURES:", "header"),
    "  DIFFUSE ST elevation",
    "  (Doesn't follow coronary territory)",
    "  PR depression in lead II - GOLD!",
    "  PR elevation in aVR",
    "  No reciprocal changes",
    "",
    ("CLINICAL CLUES:", "header"),
    "  Sharp, pleuritic chest pain",
    "  Worse with inspiration",
    "  Better sitting forward",
    "  Recent viral illness"
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(5.5))

img_path = os.path.join(IMG_DIR, "ECG-Pericarditis.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(6.3), Inches(1.5), width=Inches(6.8))

# SLIDE 18: Brugada
slide = create_content_slide("STEMI Mimic: Brugada Pattern", ACCENT_YELLOW)
content = [
    ("TYPE 1 BRUGADA:", "header"),
    "  Coved ST elevation ≥2mm in V1-V2",
    "  Followed by NEGATIVE T wave",
    "  RBBB pattern",
    "",
    ("WHY IT MATTERS:", "header"),
    "  Genetic channelopathy",
    "  Risk of sudden cardiac death",
    "  NOT a coronary occlusion",
    "  Cath lab won't find anything",
    "",
    "  Often younger patients",
    "  Family history of sudden death"
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(5.5))

img_path = os.path.join(IMG_DIR, "ECG-Brugada-Syndrome-Type-1-2.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(6.3), Inches(1.5), width=Inches(6.8))

# SLIDE 19: Hyperkalemia
slide = create_content_slide("STEMI Mimic: Hyperkalemia", ACCENT_YELLOW)
content = [
    ("PROGRESSIVE EKG CHANGES:", "header"),
    "  K+ 5.5-6.5: Peaked T waves",
    "  K+ 6.5-7.5: PR↑, P flat, QRS wide",
    "  K+ >7.5: Sine wave → VF → Asystole",
    "",
    ("WHY IT MIMICS:", "header"),
    "  Wide QRS looks like STEMI",
    "  But doesn't fit coronary territory",
    "",
    ("CLINICAL CONTEXT:", "header"),
    "  Dialysis patient?",
    "  Renal failure?",
    "  ACE-i / K-sparing diuretics?"
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(5.5))

img_path = os.path.join(IMG_DIR, "ECG-Hyperkalaemia-peaked-T-waves-serum-potassium-7.0.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(6.3), Inches(1.5), width=Inches(6.8))

# SLIDE 20: Prehospital Management - THE BIG ONE
slide = create_content_slide("Prehospital Management: 2023 Guidelines", ACCENT_GREEN)

# DO box
add_card(slide, Inches(0.3), Inches(1.5), Inches(6), Inches(4.5), "✓ DO", [
    "• 12-lead EKG within 10 min",
    "• Activate cath lab if STEMI",
    "• Aspirin 324mg CHEWED",
    "• IV access, cardiac monitor",
    "• O2 only if SpO2 <90%"
], ACCENT_GREEN)

# DON'T box
add_card(slide, Inches(6.8), Inches(1.5), Inches(6.2), Inches(4.5), "✗ DON'T", [
    "• Routine O2 in normoxic patients",
    "• Morphine as first-line",
    "• Delay transport for interventions",
    "• P2Y12 without physician order",
    "",
    "⚠️ MONA IS OUTDATED!"
], ACCENT_RED)

# Bottom mnemonic
mnem = slide.shapes.add_textbox(Inches(0.3), Inches(6.2), Inches(12.7), Inches(1))
tf = mnem.text_frame
p = tf.paragraphs[0]
p.text = "💡 NEW MNEMONIC: \"EAT\" = EKG, Aspirin, Transport"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = ACCENT_PURPLE
p.alignment = PP_ALIGN.CENTER

# SLIDE 21: Aspirin
slide = create_content_slide("Aspirin: Still the King 👑", ACCENT_GREEN)
content = [
    ("DOSE: 324mg (four 81mg tablets)", "header"),
    "",
    ("ROUTE: CHEWED, non-enteric coated", "header"),
    "",
    ("TIMING: As soon as STEMI recognized", "header"),
    "",
    ("WHY CHEWED?", "header"),
    "  Faster absorption (15-20 min vs 60 min)",
    "  Buccal absorption even if vomiting",
    "",
    ("CONTRAINDICATIONS (RARE):", "header"),
    "  True aspirin allergy",
    "  Active GI bleed",
    ("  NOT held for CABG or anticoagulation", "success")
]
add_bullet_content(slide, content, Inches(1.4))

# SLIDE 22: Nitroglycerin
slide = create_content_slide("Nitroglycerin: Use with Caution", ACCENT_YELLOW)

# Left content
content = [
    ("INDICATIONS:", "header"),
    "  Ongoing chest pain",
    "  Hypertension",
    "  Pulmonary edema",
    "",
    ("DOSE:", "header"),
    "  0.4mg SL every 5 min x3"
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(5.5))

# Right - HOLD IF box
add_card(slide, Inches(6.5), Inches(1.5), Inches(6.5), Inches(4.8), "⚠️ HOLD NITRO IF:", [
    "• SBP <90 mmHg",
    "",
    "• Viagra/Levitra within 24 hours",
    "• Cialis within 48 hours",
    "",
    "• RV infarction suspected",
    "  (inferior STEMI + hypotension)",
    "",
    "• Severe aortic stenosis"
], ACCENT_RED)

# SLIDE 23: Oxygen
slide = create_content_slide("Oxygen: Less is More", ACCENT_BLUE)
content = [
    ("2023 GUIDELINES:", "header"),
    ("  O2 ONLY if SpO2 <90%", "warning"),
    "  Target: 90-94%",
    "  AVOID hyperoxia (>96%)",
    "",
    ("WHY THE CHANGE?", "header"),
    "  AVOID trial: No benefit of routine O2",
    "  Hyperoxia may cause coronary vasoconstriction",
    "",
    ("EXCEPTIONS:", "header"),
    "  Respiratory distress",
    "  Heart failure with hypoxia",
    "  Cardiogenic shock"
]
add_bullet_content(slide, content, Inches(1.4))

# SLIDE 24: Morphine
slide = create_content_slide("Morphine: The Fallen Hero", ACCENT_YELLOW)
content = [
    ("HISTORICAL ROLE:", "header"),
    "  Pain relief, anxiety reduction, preload reduction",
    "",
    ("CURRENT EVIDENCE:", "header"),
    ("  Associated with INCREASED mortality!", "warning"),
    "  May delay P2Y12 inhibitor absorption",
    "  Can cause hypotension, respiratory depression",
    "",
    ("WHEN TO CONSIDER:", "header"),
    "  Severe pain unresponsive to nitrates",
    "  Pulmonary edema",
    "  Always have naloxone available",
    "",
    ("NOT FIRST-LINE ANYMORE", "fail")
]
add_bullet_content(slide, content, Inches(1.4))

# SLIDE 25: Cath Lab Activation
slide = create_content_slide("📱 Cath Lab Activation: You Hold the Power", ACCENT_GREEN)
content = [
    ("FIELD ACTIVATION SAVES LIVES:", "header"),
    "  Bypasses ED delay",
    "  Team mobilized before arrival",
    ("  Saves 15-20 minutes!", "success"),
    "",
    ("WHAT TO COMMUNICATE:", "header"),
    "  Patient age, sex",
    "  Symptom onset time",
    "  EKG interpretation",
    "  Vitals, mental status",
    "  Any complications",
    "  ETA"
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(6))

# Phone number callout
phone_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2), Inches(6), Inches(2.5))
phone_box.fill.solid()
phone_box.fill.fore_color.rgb = ACCENT_GREEN
phone_box.line.fill.background()

phone_text = slide.shapes.add_textbox(Inches(6.8), Inches(2.3), Inches(6), Inches(2))
tf = phone_text.text_frame
p = tf.paragraphs[0]
p.text = "📞 NORTH VISTA"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "STEMI HOTLINE"
p.font.size = Pt(24)
p.font.color.rgb = TEXT_WHITE
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "[INSERT NUMBER]"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = BG_DARK
p.alignment = PP_ALIGN.CENTER

# SLIDE 26: STEMI Treatment Algorithm (FLOWCHART)
slide = create_content_slide("🔄 STEMI Treatment Algorithm", ACCENT_BLUE)

# Flowchart boxes
steps = [
    ("CHEST PAIN\nSuspect ACS", Inches(5.5), Inches(1.5), ACCENT_YELLOW),
    ("12-Lead EKG\n(within 10 min)", Inches(5.5), Inches(2.6), ACCENT_BLUE),
    ("STEMI?", Inches(5.5), Inches(3.7), ACCENT_PURPLE),
]

# Decision branches
yes_path = [
    ("ASA 324mg\nchewed", Inches(2), Inches(4.8), ACCENT_GREEN),
    ("Activate\nCath Lab", Inches(2), Inches(5.9), ACCENT_GREEN),
]

no_path = [
    ("Consider\nMimics", Inches(9), Inches(4.8), ACCENT_YELLOW),
    ("Transport\n& Monitor", Inches(9), Inches(5.9), ACCENT_BLUE),
]

for text, x, y, color in steps:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.5), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    
    txt = slide.shapes.add_textbox(x, y + Inches(0.15), Inches(2.5), Inches(0.7))
    tf = txt.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BG_DARK if color != ACCENT_PURPLE else TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

for text, x, y, color in yes_path + no_path:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    
    txt = slide.shapes.add_textbox(x, y + Inches(0.15), Inches(2.2), Inches(0.7))
    tf = txt.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BG_DARK
    p.alignment = PP_ALIGN.CENTER

# YES/NO labels
yes_label = slide.shapes.add_textbox(Inches(3.2), Inches(4.1), Inches(1), Inches(0.5))
tf = yes_label.text_frame
p = tf.paragraphs[0]
p.text = "YES"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN

no_label = slide.shapes.add_textbox(Inches(8.5), Inches(4.1), Inches(1), Inches(0.5))
tf = no_label.text_frame
p = tf.paragraphs[0]
p.text = "NO"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT_RED

# SLIDE 27: Transport Decisions
slide = create_content_slide("Transport Decisions", ACCENT_BLUE)

add_card(slide, Inches(0.3), Inches(1.5), Inches(6), Inches(2.5), "🏥 DIRECT TO PCI CENTER IF:", [
    "• STEMI identified in field",
    "• PCI center within 60-90 minutes",
    "• Patient hemodynamically stable"
], ACCENT_GREEN)

add_card(slide, Inches(6.8), Inches(1.5), Inches(6.2), Inches(2.5), "🚑 NEAREST ED FIRST IF:", [
    "• Cardiac arrest",
    "• Cardiogenic shock",
    "• Unclear diagnosis",
    "• PCI center >120 min away"
], ACCENT_YELLOW)

add_card(slide, Inches(0.3), Inches(4.3), Inches(12.7), Inches(1.5), "🚁 CONSIDER HELICOPTER IF:", [
    "Ground transport >60 minutes to PCI center | Weather permits | Landing zone available"
], ACCENT_PURPLE)

# SLIDE 28: What Happens in Cath Lab
slide = create_content_slide("What Happens in the Cath Lab", ACCENT_BLUE)
content = [
    ("PRIMARY PCI PROCEDURE:", "header"),
    "  1. Access (radial or femoral artery)",
    "  2. Coronary angiography (find blockage)",
    "  3. Wire crosses the lesion",
    "  4. Balloon angioplasty",
    "  5. Stent deployment",
    "  6. Post-dilation and optimization",
    "",
    ("GOAL: FMC-to-Device <90 minutes", "success"),
    "",
    ("OUTCOMES:", "header"),
    "  95%+ success rate opening artery",
    "  Most patients home in 2-3 days"
]
add_bullet_content(slide, content, Inches(1.4))

# SLIDE 29: Cardiogenic Shock
slide = create_content_slide("Cardiogenic Shock: When the Pump Fails", ACCENT_RED)
content = [
    ("DEFINITION:", "header"),
    "  SBP <90 despite adequate filling",
    "  Signs of end-organ hypoperfusion",
    "  Occurs in ~5-8% of STEMI",
    ("  Mortality: 40-50%", "warning"),
    "",
    ("MANAGEMENT:", "header"),
    "  Still need PCI - primary treatment!",
    "  Vasopressors: norepinephrine preferred",
    "",
    ("MECHANICAL SUPPORT:", "header"),
    "  Impella - percutaneous LVAD",
    "  IABP - intra-aortic balloon pump",
    "  ECMO - for refractory shock"
]
add_bullet_content(slide, content, Inches(1.4))

# SLIDE 30: Medications Summary Table
slide = create_content_slide("Medications Summary", ACCENT_BLUE)

meds = [
    ("ANTIPLATELET", "Aspirin 324mg", "EMS - CHEWED!", ACCENT_GREEN),
    ("ANTIPLATELET", "P2Y12 inhibitor", "ED/Cath lab", ACCENT_BLUE),
    ("ANTICOAGULANT", "Heparin", "ED/Cath lab", ACCENT_BLUE),
    ("VASODILATOR", "Nitroglycerin", "PRN - with caution", ACCENT_YELLOW),
    ("BETA-BLOCKER", "Metoprolol", "Within 24h if stable", ACCENT_PURPLE),
    ("ACE-I/ARB", "Lisinopril", "Within 24h if stable", ACCENT_PURPLE),
    ("STATIN", "Atorvastatin 80mg", "Before discharge", ACCENT_PURPLE),
]

y = 1.5
for category, drug, timing, color in meds:
    # Category box
    cat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(y), Inches(2.8), Inches(0.65))
    cat_box.fill.solid()
    cat_box.fill.fore_color.rgb = color
    cat_box.line.fill.background()
    
    cat_text = slide.shapes.add_textbox(Inches(0.3), Inches(y + 0.12), Inches(2.8), Inches(0.5))
    tf = cat_text.text_frame
    p = tf.paragraphs[0]
    p.text = category
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BG_DARK if color not in [ACCENT_PURPLE] else TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Drug
    drug_text = slide.shapes.add_textbox(Inches(3.3), Inches(y + 0.12), Inches(4), Inches(0.5))
    tf = drug_text.text_frame
    p = tf.paragraphs[0]
    p.text = drug
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_WHITE
    
    # Timing
    time_text = slide.shapes.add_textbox(Inches(7.5), Inches(y + 0.12), Inches(5.5), Inches(0.5))
    tf = time_text.text_frame
    p = tf.paragraphs[0]
    p.text = timing
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MUTED if "EMS" not in timing else ACCENT_GREEN
    p.font.bold = "EMS" in timing
    
    y += 0.78

# YOUR JOB callout
your_job = slide.shapes.add_textbox(Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.8))
tf = your_job.text_frame
p = tf.paragraphs[0]
p.text = "💡 YOUR JOB: Aspirin. Everything else happens at the hospital."
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_YELLOW
p.alignment = PP_ALIGN.CENTER

# SLIDE 31: Cardiac Arrest with STEMI
slide = create_content_slide("Special Scenario: Cardiac Arrest with STEMI", ACCENT_RED)
content = [
    ("IF STEMI ON POST-ROSC EKG:", "header"),
    "  Immediate cath lab activation",
    "  Transport to PCI center",
    "  Consider targeted temperature management",
    "  May need mechanical support",
    "",
    ("IF NO STEMI BUT SUSPECTED CARDIAC:", "header"),
    "  May still benefit from angiography",
    "  ~30% of OHCA have culprit lesion",
    "  Discuss with receiving facility"
]
add_bullet_content(slide, content, Inches(1.4))

# SLIDE 32: Cocaine-Associated MI
slide = create_content_slide("Special Scenario: Cocaine-Associated MI", ACCENT_PURPLE)
content = [
    ("MECHANISM:", "header"),
    "  Coronary vasospasm",
    "  Accelerated atherosclerosis",
    "  Prothrombotic state",
    "  Increased O2 demand",
    "",
    ("MANAGEMENT DIFFERENCES:", "header"),
    "  Benzos first-line for agitation",
    ("  ⚠️ AVOID BETA-BLOCKERS!", "warning"),
    "  (unopposed alpha effect)",
    "  Nitrates are helpful",
    "  PCI if true STEMI with thrombus"
]
add_bullet_content(slide, content, Inches(1.4))

# SLIDE 33-35: Case Studies
# Case 1
slide = create_content_slide("Case Study 1: 62-year-old Male", ACCENT_BLUE)
content = [
    ("PRESENTATION:", "header"),
    "  Crushing chest pressure x 45 min",
    "  Hx: HTN, DM, former smoker",
    "  Vitals: HR 88, BP 156/92",
    "  SpO2 97% RA",
    "",
    ("QUESTIONS:", "header"),
    "  1. What type of STEMI?",
    "  2. What artery is occluded?",
    "  3. Immediate actions?"
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(5.5))

img_path = os.path.join(IMG_DIR, "ECG-Extensive-anterior-MI-tombstoning-pattern.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(6.3), Inches(1.5), width=Inches(6.8))

# Case 2
slide = create_content_slide("Case Study 2: 71-year-old Female", ACCENT_YELLOW)
content = [
    ("PRESENTATION:", "header"),
    "  Chest discomfort, nausea x 2 hrs",
    "  Hx: HTN, CKD stage 3",
    "  Vitals: HR 48, BP 88/54",
    "  SpO2 94% RA",
    "",
    ("QUESTIONS:", "header"),
    "  1. What type of STEMI?",
    "  2. Why bradycardic/hypotensive?",
    "  3. What to do (and NOT do)?"
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(5.5))

img_path = os.path.join(IMG_DIR, "ECG-Inferior-AMI-STEMI-Massive-tombstones.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(6.3), Inches(1.5), width=Inches(6.8))

# Case 3
slide = create_content_slide("Case Study 3: 28-year-old Male", ACCENT_PURPLE)
content = [
    ("PRESENTATION:", "header"),
    "  Sharp chest pain, worse breathing",
    "  Hx: URI 1 week ago",
    "  Vitals: HR 102, BP 128/78",
    "  SpO2 99% RA",
    "",
    ("QUESTIONS:", "header"),
    "  1. Is this STEMI?",
    "  2. Key differentiating features?",
    "  3. Activate cath lab?"
]
add_bullet_content(slide, content, Inches(1.4), Inches(0.3), Inches(5.5))

img_path = os.path.join(IMG_DIR, "ECG-Pericarditis.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(6.3), Inches(1.5), width=Inches(6.8))

# SLIDE 36: Key Take-Home Points
slide = create_content_slide("🎯 Key Take-Home Points", ACCENT_GREEN)

points = [
    ("1", "TIME = MUSCLE", "Every minute counts"),
    ("2", "EKG EARLY", "Within 10 minutes"),
    ("3", "ACTIVATE EARLY", "Saves 15-20 min"),
    ("4", "KNOW MIMICS", "But when in doubt, activate"),
    ("5", "ASPIRIN IS KING", "324mg chewed"),
    ("6", "CHECK V4R", "On inferior STEMIs"),
    ("7", "O2 IF NEEDED", "Only if SpO2 <90%"),
    ("8", "LESS MORPHINE", "Use judiciously"),
]

for i, (num, title, detail) in enumerate(points):
    row = i // 2
    col = i % 2
    x = 0.3 + (col * 6.5)
    y = 1.5 + (row * 1.4)
    
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_GREEN
    circle.line.fill.background()
    
    num_text = slide.shapes.add_textbox(Inches(x), Inches(y + 0.05), Inches(0.5), Inches(0.4))
    tf = num_text.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BG_DARK
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_text = slide.shapes.add_textbox(Inches(x + 0.6), Inches(y), Inches(5.5), Inches(0.4))
    tf = title_text.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    # Detail
    detail_text = slide.shapes.add_textbox(Inches(x + 0.6), Inches(y + 0.4), Inches(5.5), Inches(0.4))
    tf = detail_text.text_frame
    p = tf.paragraphs[0]
    p.text = detail
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MUTED

# SLIDE 37: Quick Reference Checklist
slide = create_content_slide("📋 STEMI Field Checklist", ACCENT_GREEN)

checklist = [
    "☐ 12-lead EKG (within 10 min)",
    "☐ Identify STEMI pattern",
    "☐ Call STEMI hotline",
    "☐ Aspirin 324mg (chewed)",
    "☐ O2 only if SpO2 <90%",
    "☐ IV access",
    "☐ Nitro (if no contraindication)",
    "☐ Check V4R if inferior STEMI",
    "☐ Transport CODE 3",
]

y = 1.5
for item in checklist:
    item_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(6), Inches(0.5))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_WHITE
    y += 0.55

# HOLD NITRO box
add_card(slide, Inches(7), Inches(1.5), Inches(6), Inches(3), "⚠️ HOLD NITRO IF:", [
    "• SBP <90",
    "• RV infarct suspected",
    "• PDE-5 inhibitor (24-48h)"
], ACCENT_RED)

# V4R reminder
add_card(slide, Inches(7), Inches(4.8), Inches(6), Inches(1.8), "💡 REMEMBER V4R", [
    "On ALL inferior STEMIs!",
    "RV infarct changes everything"
], ACCENT_YELLOW)

# SLIDE 38: North Vista Protocol
slide = create_content_slide("North Vista STEMI Protocol", ACCENT_GREEN)

# Big phone number
phone_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(1.5), Inches(7.333), Inches(1.8))
phone_box.fill.solid()
phone_box.fill.fore_color.rgb = ACCENT_GREEN
phone_box.line.fill.background()

phone_text = slide.shapes.add_textbox(Inches(3), Inches(1.7), Inches(7.333), Inches(1.5))
tf = phone_text.text_frame
p = tf.paragraphs[0]
p.text = "📞 STEMI HOTLINE"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = BG_DARK
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "[INSERT NUMBER]"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.alignment = PP_ALIGN.CENTER

# Two columns
add_card(slide, Inches(0.3), Inches(3.6), Inches(6), Inches(3.5), "📤 WHAT WE NEED FROM YOU:", [
    "• Call early (even on scene)",
    "• Patient: age, sex, symptom onset",
    "• EKG interpretation",
    "• Vitals, mental status",
    "• Any complications",
    "• ETA"
], ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(3.6), Inches(6.2), Inches(3.5), "✓ WHAT WE'LL HAVE READY:", [
    "• Cath lab team gowned",
    "• Anesthesia on standby",
    "• Direct gurney-to-table pathway",
    "",
    "We're waiting for YOU!"
], ACCENT_GREEN)

# SLIDE 39: Questions
create_title_slide(
    "Questions?",
    [
        "",
        "Dr. Peter Bleszynski, MD, FACC, FSCAI",
        "Director of Cardiology, North Vista Hospital",
        "",
        "📞 STEMI Hotline: [INSERT NUMBER]"
    ]
)

# SLIDE 40: Resources
slide = create_content_slide("Resources for Further Learning", ACCENT_BLUE)
content = [
    ("GUIDELINES:", "header"),
    "  2023 AHA/ACC/ACEP/NAEMSP/SCAI Guideline for ACS",
    "  2021 ACC/AHA/SCAI Guideline for Coronary Revascularization",
    "",
    ("EKG PRACTICE:", "header"),
    "  Life in the Fast Lane (litfl.com)",
    "  Dr. Smith's ECG Blog",
    "  ECG Wave-Maven",
    "",
    ("APPS:", "header"),
    "  AHA ECG Course",
    "  Medscape"
]
add_bullet_content(slide, content, Inches(1.4))

# SLIDE 41: Thank You
create_title_slide(
    "Thank You",
    [
        "You save lives every day.",
        "",
        "🫀",
        "",
        "Questions? Contact Dr. Bleszynski"
    ]
)

# ============== SAVE ==============
output_path = "/Users/vi/.openclaw/workspace/projects/stemi-presentation/STEMI_EMS_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"🎨 Features included:")
print(f"   - EKG images integrated")
print(f"   - Treatment algorithm flowchart")
print(f"   - Timeline infographic")
print(f"   - Memorable mnemonics")
print(f"   - Visual cards and color-coding")
print(f"   - Quick reference checklist")
